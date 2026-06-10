"""Extract plain text from uploaded course materials."""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path


@dataclass(frozen=True)
class PageText:
    page_number: int
    text: str


def extract_text(path: str | Path, filename: str | None = None) -> str:
    """Flat text extraction (legacy fallback)."""
    pages = extract_pages(path, filename)
    return "\n\n".join(p.text for p in pages if p.text.strip())


def extract_pages(path: str | Path, filename: str | None = None) -> list[PageText]:
    p = Path(path)
    if not p.is_file():
        raise FileNotFoundError(f"Material file not found: {p}")

    name = (filename or p.name).lower()
    if name.endswith(".pdf"):
        return _extract_pdf_pages(p)
    if name.endswith(".pptx"):
        return _extract_pptx_pages(p)
    if name.endswith(".txt") or name.endswith(".md"):
        text = p.read_text(encoding="utf-8", errors="replace").strip()
        return [PageText(page_number=1, text=text)] if text else []
    raise ValueError(f"Unsupported file type for text extraction: {name}")


def _extract_pdf_pages(path: Path) -> list[PageText]:
    import fitz  # PyMuPDF

    doc = fitz.open(path)
    pages: list[PageText] = []
    try:
        for i, page in enumerate(doc, start=1):
            text = page.get_text().strip()
            if text:
                pages.append(PageText(page_number=i, text=text))
    finally:
        doc.close()
    return pages


def _extract_pptx_pages(path: Path) -> list[PageText]:
    from pptx import Presentation

    prs = Presentation(str(path))
    pages: list[PageText] = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_bits: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slide_bits.append(shape.text.strip())
        if slide_bits:
            pages.append(PageText(page_number=i, text="\n".join(slide_bits)))
    return pages
